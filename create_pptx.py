from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

base = "/Users/thom/Documents/Ausbildung/Schule/Website trash Island"
logo = os.path.join(base, "Logo_Trash Island_Electrisize_weiss.png")
img3 = os.path.join(base, "3.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(11, 11, 13)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(190, 190, 200)
ACCENT = RGBColor(48, 253, 71)
ORANGE = RGBColor(255, 122, 0)


def set_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()


def add_text(slide, text, x, y, w, h, size, color, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box

# Slide 1
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1)
add_text(s1, "Website Design - Trash Island", 0.9, 0.9, 11.6, 1.0, 46, WHITE, True)
add_text(s1, "Apple-Style UI: Glas, klare Typografie, starke Akzente", 0.9, 2.0, 9.8, 1.0, 20, MUTED)
if os.path.exists(logo):
    s1.shapes.add_picture(logo, Inches(0.9), Inches(5.5), height=Inches(1.4))
line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.2), Inches(4.2), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Slide 2
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2)
add_text(s2, "Design-Highlights", 0.9, 0.7, 8.0, 1.0, 40, WHITE, True)
add_text(s2, "- Blur-Navigation mit weichen Schatten\n- Gruen-Orange Akzente\n- Animierte Buttons und Uebergaenge\n- Fisheye Typografie", 0.9, 1.8, 6.4, 3.6, 20, WHITE)
if os.path.exists(img3):
    s2.shapes.add_picture(img3, Inches(7.6), Inches(1.6), height=Inches(4.9))
chip = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.9), Inches(3.5), Inches(0.5))
chip.fill.solid()
chip.fill.fore_color.rgb = RGBColor(20, 20, 26)
chip.line.color.rgb = ACCENT
ctf = chip.text_frame
ctf.text = "Apple-Style Layout"
ctf.paragraphs[0].font.size = Pt(14)
ctf.paragraphs[0].font.color.rgb = ACCENT

# Slide 3
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3)
add_text(s3, "Rechtliche Hinweise - Modul", 0.9, 0.7, 10.0, 1.0, 40, WHITE, True)
add_text(s3, "Kompaktes Fenster mit Toolbar und klarem Fokus", 0.9, 1.8, 9.5, 0.7, 18, MUTED)
card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.6), Inches(7.8), Inches(3.6))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(18, 18, 24)
card.line.color.rgb = RGBColor(50, 50, 60)
bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.6), Inches(7.8), Inches(0.5))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(28, 28, 36)
bar.line.fill.background()

for i, c in enumerate([RGBColor(255,95,87), RGBColor(255,189,46), RGBColor(40,200,64)]):
    d = s3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0 + i*0.25), Inches(2.72), Inches(0.18), Inches(0.18))
    d.fill.solid()
    d.fill.fore_color.rgb = c
    d.line.fill.background()

add_text(s3, "- Hygiene und Abfallrichtlinien (LMHV)\n- EU-Lebensmittelrecht (EG 852/2004)\n- KrWG: Vermeidung, Verwertung, Entsorgung", 1.2, 3.2, 7.1, 2.5, 18, WHITE)
if os.path.exists(logo):
    s3.shapes.add_picture(logo, Inches(9.3), Inches(2.8), height=Inches(1.6))
line = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.4), Inches(4.5), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = ORANGE
line.line.fill.background()

out_path = os.path.join(base, "Website-Design-Trash-Island.pptx")
prs.save(out_path)
print(out_path)
